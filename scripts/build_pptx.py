#!/usr/bin/env python3
"""V3 — codex's rescue moves applied:
   1) Type scale tamed (display 56pt down from 72; body 22pt; controlled bold).
   2) Auto-size SHAPE_TO_FIT_TEXT on variable-content frames (no silent clipping).
   3) Open composition for application/prayer (no boxed cards).
   + Image per day (hero right half).
   + Korean East Asian font handling (latin/ea/cs).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

ROOT = Path("/Users/csh/Projects/claude-code-md-explainer/devotion-10")
IMG = ROOT / "img"

# === Tokens — Anthropic claude.com palette ===
C = {
    "primary":      RGBColor(0xCC, 0x78, 0x5C),
    "canvas":       RGBColor(0xFA, 0xF9, 0xF5),
    "card":         RGBColor(0xEF, 0xE9, 0xDE),
    "dark":         RGBColor(0x18, 0x17, 0x15),
    "ink":          RGBColor(0x14, 0x14, 0x13),
    "body":         RGBColor(0x3D, 0x3D, 0x3A),
    "muted":        RGBColor(0x6C, 0x6A, 0x64),
    "on_dark":      RGBColor(0xFA, 0xF9, 0xF5),
    "on_dark_soft": RGBColor(0xA0, 0x9D, 0x96),
    "hairline":     RGBColor(0xE6, 0xDF, 0xD8),
}

FONT_DISPLAY = "NanumMyeongjo"
FONT_SANS = "SUIT"

# === Korean font handling — set latin + ea + cs ===
def set_korean_font(run, font_name, bold=False):
    rPr = run._r.get_or_add_rPr()
    if bold:
        rPr.set("b", "1")
    for tag in ["a:latin", "a:ea", "a:cs"]:
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    for tag in ["a:latin", "a:ea", "a:cs"]:
        e = etree.SubElement(rPr, qn(tag))
        e.set("typeface", font_name)

def add_text(slide, x, y, w, h, text, *, font=FONT_SANS, size=18, bold=False,
             color=None, align="left", valign="top", line_spacing=1.5,
             auto_size=False, char_spacing=0):
    if color is None:
        color = C["ink"]
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    if auto_size:
        # Box grows downward to fit text — no silent clipping.
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    else:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        set_korean_font(run, font, bold=bold)
        if char_spacing:
            run._r.get_or_add_rPr().set("spc", str(char_spacing))
    return box

def add_rect(slide, x, y, w, h, fill, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape

# === Layout — 16:9 at 13.333 × 7.5 inch ===
W, H = 13.333, 7.5
PAD = 0.85
CW = W - 2 * PAD

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]
TOTAL = 16

def bg_canvas(slide, color):
    shape = add_rect(slide, 0, 0, W, H, color)
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)

def page_footer(slide, n, dark=False):
    cap_color = C["on_dark_soft"] if dark else C["muted"]
    add_text(slide, PAD, H - 0.55, 7, 0.4,
             "한 가족이 되신 예수님 · 5일 묵상",
             font=FONT_SANS, size=16, bold=True, color=cap_color, char_spacing=300)
    add_text(slide, W - 2.0, H - 0.55, 1.8, 0.4,
             f"{n:02d} / {TOTAL:02d}",
             font=FONT_SANS, size=16, bold=True, color=cap_color, align="right", char_spacing=300)

def label(slide, text, dark=False, y=PAD - 0.05):
    color = C["on_dark"] if dark else C["muted"]
    add_text(slide, PAD, y, CW, 0.45, text,
             font=FONT_SANS, size=22, bold=True, color=color, char_spacing=400)

# ============================================================
# Slide 1: Series intro
# ============================================================
s = prs.slides.add_slide(BLANK)
bg_canvas(s, C["canvas"])
label(s, "사도신경 강해 10 · 죽으시고, 장사된 지")
add_text(s, PAD, 1.4, CW, 2.0,
         "한 가족이 되신 예수님,\n다섯 날의 묵상.",
         font=FONT_DISPLAY, size=56, bold=True, color=C["ink"], line_spacing=1.1)
add_text(s, PAD, 3.5, CW * 0.85, 0.7,
         "매일 한 편씩 — 본문 한 구절, 핵심 한 문장, 짧은 묵상글과 기도.",
         font=FONT_SANS, size=24, bold=True, color=C["body"], auto_size=True)

# 5-day grid
tile_w = (CW - 0.6) / 5
tile_h = 2.5
tile_y = 4.4
tiles = [
    ("DAY 01", "고후 8:9", "한 가족이 되신\n예수님."),
    ("DAY 02", "고후 8:9", "가난과 부요의\n교환."),
    ("DAY 03", "갈 3:13-14", "저주를 짊어지신\n예수님."),
    ("DAY 04", "고후 5:21", "의의 전가, 새로운\n신분."),
    ("DAY 05", "롬 8:38-39", "죽음조차 끊지\n못하는 사랑."),
]
for i, (num, scrip, ttl) in enumerate(tiles):
    x = PAD + i * (tile_w + 0.15)
    # Subtle hairline border, no fill — minimal
    rect = add_rect(s, x, tile_y, tile_w, tile_h, C["canvas"], line_color=C["hairline"])
    add_text(s, x + 0.25, tile_y + 0.3, tile_w - 0.5, 0.4,
             num, font=FONT_SANS, size=18, bold=True, color=C["primary"], char_spacing=600)
    add_text(s, x + 0.25, tile_y + 0.75, tile_w - 0.5, 0.4,
             scrip, font=FONT_SANS, size=16, bold=True, color=C["muted"], char_spacing=400)
    add_text(s, x + 0.25, tile_y + tile_h - 1.4, tile_w - 0.5, 1.3,
             ttl, font=FONT_DISPLAY, size=24, bold=True, color=C["ink"],
             line_spacing=1.18, valign="bottom")
page_footer(s, 1)

# ============================================================
# Day-builder
# ============================================================
def build_day(num, label_txt, ttl1, ttl2, sub, scrip, scrip_cite,
              key1, key2, med_p1, med_p2, app_h, q1, q2, prayer,
              hero_bg="canvas", img_name=None):
    base_idx = 1 + (num - 1) * 3 + 1  # 2, 5, 8, 11, 14

    # === Slide A: Hero with image (left text, right image) ===
    s = prs.slides.add_slide(BLANK)
    is_dark = hero_bg == "dark"
    bg_color = {"canvas": C["canvas"], "card": C["card"], "dark": C["dark"]}[hero_bg]
    bg_canvas(s, bg_color)

    text_w = 6.5  # left column width
    label(s, f"DAY 0{num} · {label_txt}", dark=is_dark)
    # Display title
    add_text(s, PAD, 1.4, text_w, 2.5,
             f"{ttl1}\n{ttl2}",
             font=FONT_DISPLAY, size=58, bold=True,
             color=(C["on_dark"] if is_dark else C["ink"]),
             line_spacing=1.1, auto_size=True)
    # Subhead
    add_text(s, PAD, 4.1, text_w, 1.1, sub,
             font=FONT_SANS, size=24, bold=True,
             color=(C["on_dark_soft"] if is_dark else C["body"]),
             line_spacing=1.45, auto_size=True)
    # Scripture: open with coral left border
    add_rect(s, PAD, 5.35, 0.06, 1.6, C["primary"])
    add_text(s, PAD + 0.22, 5.30, text_w - 0.22, 1.6,
             f'"{scrip}"',
             font=FONT_DISPLAY, size=28, bold=True,
             color=(C["on_dark"] if is_dark else C["ink"]),
             line_spacing=1.34, auto_size=True)
    add_text(s, PAD + 0.22, 6.75, text_w - 0.22, 0.4,
             scrip_cite,
             font=FONT_SANS, size=21, bold=True,
             color=(C["on_dark_soft"] if is_dark else C["muted"]),
             char_spacing=400)

    # Image on right
    if img_name and (IMG / img_name).exists():
        img_x = 7.7
        img_y = 1.2
        img_size = 5.3
        s.shapes.add_picture(str(IMG / img_name),
            Inches(img_x), Inches(img_y), Inches(img_size), Inches(img_size))

    page_footer(s, base_idx, dark=is_dark)

    # === Slide B: Meditation (clean cream, key sentence + body) ===
    s = prs.slides.add_slide(BLANK)
    bg_canvas(s, C["canvas"])
    label(s, f"DAY 0{num} · 묵상")
    # Key pull quote — wide enough that single line never wraps mid-word
    add_text(s, PAD, 1.4, CW, 2.2,
             f"{key1}\n{key2}",
             font=FONT_DISPLAY, size=46, bold=True, color=C["ink"],
             line_spacing=1.2, auto_size=True)
    # Meditation paragraphs (auto-size guards against overflow)
    add_text(s, PAD, 3.9, CW * 0.92, 1.5, med_p1,
             font=FONT_SANS, size=24, bold=True, color=C["body"],
             line_spacing=1.65, auto_size=True)
    add_text(s, PAD, 5.6, CW * 0.92, 1.4, med_p2,
             font=FONT_SANS, size=24, bold=True, color=C["body"],
             line_spacing=1.65, auto_size=True)
    page_footer(s, base_idx + 1)

    # === Slide C: Application + Prayer (open composition, no boxes) ===
    s = prs.slides.add_slide(BLANK)
    bg_canvas(s, C["card"])
    label(s, f"DAY 0{num} · 적용 + 기도")
    # App headline
    add_text(s, PAD, 1.4, CW, 1.4, app_h,
             font=FONT_DISPLAY, size=38, bold=True, color=C["ink"],
             line_spacing=1.25, auto_size=True)
    # Two questions — open type, just numbers + text, no boxes
    add_text(s, PAD, 2.95, 0.8, 0.4,
             "01", font=FONT_SANS, size=18, bold=True, color=C["primary"], char_spacing=400)
    add_text(s, PAD + 0.8, 2.9, CW - 1.0, 1.0, q1,
             font=FONT_SANS, size=22, bold=True, color=C["ink"],
             line_spacing=1.55, auto_size=True)
    add_text(s, PAD, 4.25, 0.8, 0.4,
             "02", font=FONT_SANS, size=18, bold=True, color=C["primary"], char_spacing=400)
    add_text(s, PAD + 0.8, 4.20, CW - 1.0, 1.0, q2,
             font=FONT_SANS, size=22, bold=True, color=C["ink"],
             line_spacing=1.55, auto_size=True)
    # Coral hairline divider
    add_rect(s, PAD, 5.6, 1.0, 0.04, C["primary"])
    # Prayer
    add_text(s, PAD, 5.85, CW * 0.92, 1.4, prayer,
             font=FONT_DISPLAY, size=26, bold=True, color=C["ink"],
             line_spacing=1.55, auto_size=True)
    page_footer(s, base_idx + 2)

# ============================================================
# Days 1-5
# ============================================================
build_day(1, "한 가족 되심",
    "한 가족이", "되신 예수님.",
    "예수님은 우리를 도와주신 게 아니라, 우리와 한 가족이 되셨습니다.",
    "그리스도께서는 부요하나, 여러분을 위해서 가난하게 되셨습니다.",
    "고린도후서 8:9",
    "가족이 된다는 건,", "모든 것을 함께 나눈다는 뜻입니다.",
    "부모님 카드 속 돈은 부모님이 번 것이지만, 가족이라는 이유 하나로 내가 쓸 수 있습니다. 가족이 된다는 건 그런 겁니다.",
    "예수님이 우리에게 오신 방식이 그렇습니다. 우리의 모든 것 — 가난도, 저주도, 죄도 — 자기 것으로 가져가시고, 그분의 부요와 의로움을 우리에게 선물로 주셨습니다.",
    "한 가지만 골라, 오늘 마음에 두기.",
    "오늘 내가 예수님께 \"가족이라서\" 가져가시도록 맡길 수 있는 한 가지는 무엇입니까?",
    "\"가족이 되어 주셨다\"는 표현을 들었을 때, 가장 먼저 떠오르는 장면은 무엇입니까?",
    "예수님, 저를 돕는 분이 아니라 한 가족이 되어 주셔서 감사합니다.\n오늘 제 짐을 혼자 지지 않게 하소서.",
    hero_bg="canvas", img_name="day1.png")

build_day(2, "부요와 가난의 교환",
    "가난과", "부요의 교환.",
    "그분이 가난해지셨기에, 나는 부요해졌습니다.",
    "그의 가난으로 여러분을 부요하게 하시려는 것입니다.",
    "고린도후서 8:9",
    "도와주신 게 아니라,", "내 자리에 직접 내려오셨습니다.",
    "예수님은 후원자가 아니라 우리 자리에 직접 내려오신 분입니다. 그분이 가난하셨기에 우리는 부요해졌습니다.",
    "가난은 단지 돈이 없는 것만 아닙니다. 사랑받지 못함, 인정받지 못함, 안전감의 부재 — 이 모든 가난을 그분이 먼저 지나가셨습니다.",
    "내 안의 \"가난\"을 한 단어로 적어보기.",
    "내가 지금 느끼는 \"가난\"은 무엇입니까? 그것이 예수님이 이미 지나가신 자리임을 떠올려 보세요.",
    "오늘 누군가에게 \"내가 가진 부요\"를 작은 형태로라도 나눌 수 있는 한 가지는 무엇입니까?",
    "예수님, 저를 위해 가난해지신 그 사랑을 묵상합니다.\n오늘도 누군가에게 작은 부요를 흘려보내게 하소서.",
    hero_bg="canvas", img_name="day2.png")

build_day(3, "저주의 자리",
    "저주를", "짊어지신 예수님.",
    "가장 큰 저주의 자리에 그분이 가셨기에, 가장 큰 축복이 내게 흘러옵니다.",
    "그리스도께서 우리를 위하여 저주를 받은 사람이 되심으로써, 우리를 율법의 저주에서 속량해 주셨습니다.",
    "갈라디아서 3:13",
    "예수님은 두려운 자리에", "먼저 가신 분입니다.",
    "십자가는 단순한 처형 도구가 아니었습니다. 당시 사람들에게 그것은 \"하나님께 버림받았다\"는 가장 큰 모욕이자 저주였습니다.",
    "그 자리에 예수님이 가셨습니다. 자기 잘못이 하나도 없으셨는데도 — 우리가 받아야 할 저주를 자기 등에 지셨습니다.",
    "내 안의 \"버림받음\"의 그림자를 마주하기.",
    "내가 두려워하는 \"버림받음\"의 형태는 무엇입니까? 그 자리에도 예수님이 먼저 가셨다는 사실이 위로가 됩니까?",
    "오늘 내가 누군가의 \"저주\" — 외로움이나 수치 — 를 함께 짊어질 수 있다면 무엇입니까?",
    "예수님, 가장 깊은 어둠의 자리에 먼저 가셔서 감사합니다.\n그 발자국을 따라 한 걸음 내딛게 하소서.",
    hero_bg="dark", img_name="day3.png")

build_day(4, "새 신분",
    "의의 전가,", "새로운 신분.",
    "내가 의로워진 게 아니라, 그분의 의로움이 내게 입혀졌습니다.",
    "그것은 우리가 그리스도 안에서 하나님의 의가 되게 하시려는 것입니다.",
    "고린도후서 5:21",
    "하나님은 내가 한 일이 아니라,", "예수님이 하신 일을 보십니다.",
    "\"내가 노력해서 깨끗해지면 하나님이 받아주시겠지\" — 우리는 자주 이렇게 생각합니다. 그러나 본문은 정반대로 말합니다.",
    "죄가 전혀 없으신 분이 우리 자리에 죄인으로 서셨고, 그래서 죄 많은 우리가 그분 안에서 의로운 사람이 되었습니다. 이것은 거래가 아니라 가족의 교환입니다.",
    "\"내 의로움\"을 내려놓는 한 순간.",
    "\"내 의로움\"이 아니라 \"그분의 의로움\" 안에서 산다면, 무엇이 가장 가벼워질까요?",
    "누군가를 그의 행동이 아니라 \"예수님 안의 신분\"으로 봐주는 한 번의 시선을 가져 보세요.",
    "예수님, 제가 한 일이 아니라 당신이 하신 일을 의지하게 하소서.\n당신 안의 새 신분으로 살게 하소서.",
    hero_bg="canvas", img_name="day4.png")

build_day(5, "끊어지지 않는 사랑",
    "죽음조차", "끊지 못하는 사랑.",
    "그분이 죽음의 자리까지 가셨기에, 그곳도 더는 우리를 끊지 못합니다.",
    "그 어떤 피조물도, 우리를 우리 주 예수 그리스도 안에 있는 하나님의 사랑에서 끊을 수 없습니다.",
    "로마서 8:38-39",
    "예수님이 먼저 가신 자리는,", "더 이상 우리를 끊지 못합니다.",
    "예수님의 죽음은 연극이 아니었습니다. 진짜로 심장이 멈추고, 무덤에 묻히신 사건입니다. 우리와 완전히 하나가 되시기 위해서였습니다.",
    "외로움, 실패, 불안, 두려움 — 우리가 가는 모든 자리에 예수님은 먼저 가셨고, 거기서 우리에게 손을 내미십니다.",
    "한 주의 묵상에서 한 가지를 챙겨가기.",
    "이번 한 주 묵상에서 \"예수님이 먼저 가신 자리\"라는 표현이 가장 깊이 와닿은 순간은 무엇입니까?",
    "다음 한 주, 어떤 한 가지 영역에서 이 진리를 살아내고 싶습니까?",
    "예수님, 죽음의 자리까지 가신 그 사랑이 제 매일의 자리가 되게 하소서.\n이 한 주 묵상이 끝이 아니라, 일상의 시작이 되게 하소서.",
    hero_bg="canvas", img_name="day5.png")

out_path = ROOT / "devotional_v3.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slide count: {len(prs.slides)}")
